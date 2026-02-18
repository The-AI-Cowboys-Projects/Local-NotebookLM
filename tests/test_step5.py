"""Tests for step5 — infographic generation."""

import json
import pickle
import pytest
from unittest.mock import MagicMock, patch

from local_notebooklm.steps.step5 import (
    InfographicError,
    extract_structured_data,
    load_transcript_text,
    render_infographic_html,
    render_png,
    step5,
)
from local_notebooklm.steps.step5_pptx import (
    render_infographic_pptx,
    render_video,
)
from local_notebooklm.steps.step5_charts import (
    generate_all_charts,
    generate_conversation_flow_chart,
    generate_speaker_distribution_chart,
    generate_topic_importance_chart,
)


# ---------------------------------------------------------------------------
# Sample structured data for rendering tests
# ---------------------------------------------------------------------------

SAMPLE_DATA = {
    "title": "AI Revolution",
    "summary": "A deep dive into modern AI capabilities.",
    "topics": [
        {"name": "Neural Networks", "description": "How NNs work", "importance": 5},
        {"name": "Ethics", "description": "AI ethics concerns", "importance": 3},
    ],
    "key_takeaways": [
        "AI is transforming industries",
        "Ethics remain critical",
        "Open-source models are rising",
    ],
    "notable_quotes": [
        {"speaker": "Speaker 1", "quote": "This changes everything."},
        {"speaker": "Speaker 2", "quote": "We need guardrails."},
    ],
    "speakers": [
        {"label": "Speaker 1", "role": "Host", "line_count": 30},
        {"label": "Speaker 2", "role": "Co-Host", "line_count": 25},
    ],
    "conversation_flow": [
        {"speaker": "Speaker 1", "topic": "Introduction"},
        {"speaker": "Speaker 2", "topic": "Background"},
        {"speaker": "Speaker 1", "topic": "Deep dive"},
        {"speaker": "Speaker 2", "topic": "Implications"},
        {"speaker": "Speaker 1", "topic": "Conclusion"},
    ],
}


# ---------------------------------------------------------------------------
# TestLoadTranscriptText
# ---------------------------------------------------------------------------

class TestLoadTranscriptText:
    def test_loads_txt_file(self, tmp_path):
        (tmp_path / "podcast_ready_data.txt").write_text("Hello world", encoding="utf-8")
        result = load_transcript_text(str(tmp_path))
        assert result == "Hello world"

    def test_falls_back_to_pkl(self, tmp_path):
        data = [("Speaker 1", "Hello"), ("Speaker 2", "Hi")]
        with open(tmp_path / "podcast_ready_data.pkl", "wb") as f:
            pickle.dump(data, f)
        result = load_transcript_text(str(tmp_path))
        assert "Speaker 1: Hello" in result
        assert "Speaker 2: Hi" in result

    def test_txt_takes_priority_over_pkl(self, tmp_path):
        (tmp_path / "podcast_ready_data.txt").write_text("from txt", encoding="utf-8")
        with open(tmp_path / "podcast_ready_data.pkl", "wb") as f:
            pickle.dump("from pkl", f)
        result = load_transcript_text(str(tmp_path))
        assert result == "from txt"

    def test_missing_files_raises(self, tmp_path):
        with pytest.raises(InfographicError, match="No transcript found"):
            load_transcript_text(str(tmp_path))

    def test_pkl_string_format(self, tmp_path):
        with open(tmp_path / "podcast_ready_data.pkl", "wb") as f:
            pickle.dump("raw string transcript", f)
        result = load_transcript_text(str(tmp_path))
        assert result == "raw string transcript"


# ---------------------------------------------------------------------------
# TestExtractStructuredData
# ---------------------------------------------------------------------------

class TestExtractStructuredData:
    def _make_config(self):
        return {
            "Big-Text-Model": {"model": "test-model"},
            "Step5": {"max_tokens": 4096, "temperature": 0.4},
        }

    @patch("local_notebooklm.steps.step5.generate_text")
    def test_valid_json_extraction(self, mock_gen):
        mock_gen.return_value = json.dumps(SAMPLE_DATA)
        result = extract_structured_data(MagicMock(), self._make_config(), "transcript")
        assert result["title"] == "AI Revolution"
        assert len(result["topics"]) == 2

    @patch("local_notebooklm.steps.step5.generate_text")
    def test_strips_markdown_fences(self, mock_gen):
        mock_gen.return_value = "```json\n" + json.dumps(SAMPLE_DATA) + "\n```"
        result = extract_structured_data(MagicMock(), self._make_config(), "transcript")
        assert result["title"] == "AI Revolution"

    @patch("local_notebooklm.steps.step5.generate_text")
    def test_invalid_json_raises(self, mock_gen):
        mock_gen.return_value = "this is not json at all"
        with pytest.raises(InfographicError, match="invalid JSON"):
            extract_structured_data(MagicMock(), self._make_config(), "transcript")

    @patch("local_notebooklm.steps.step5.generate_text")
    def test_missing_keys_raises(self, mock_gen):
        mock_gen.return_value = json.dumps({"title": "only title"})
        with pytest.raises(InfographicError, match="missing keys"):
            extract_structured_data(MagicMock(), self._make_config(), "transcript")


# ---------------------------------------------------------------------------
# TestRenderInfographicHtml
# ---------------------------------------------------------------------------

class TestRenderInfographicHtml:
    def test_returns_html_string(self):
        result = render_infographic_html(SAMPLE_DATA)
        assert "<!DOCTYPE html>" in result
        assert "AI Revolution" in result
        assert "</html>" in result

    def test_no_external_urls(self):
        result = render_infographic_html(SAMPLE_DATA)
        assert "http://" not in result
        assert "https://" not in result

    def test_xss_escaping(self):
        xss_data = {
            **SAMPLE_DATA,
            "title": '<script>alert("xss")</script>',
            "summary": '"><img src=x onerror=alert(1)>',
        }
        result = render_infographic_html(xss_data)
        assert "<script>" not in result
        assert "&lt;script&gt;" in result
        # The < and > are escaped so no actual HTML tag is injected
        assert "<img " not in result
        assert "&lt;img " in result

    def test_all_sections_present(self):
        result = render_infographic_html(SAMPLE_DATA)
        assert "// Topics" in result
        assert "// Key Takeaways" in result
        assert "// Notable Quotes" in result
        assert "// Speakers" in result
        assert "// Conversation Flow" in result
        assert "Generated by Local-NotebookLM" in result


# ---------------------------------------------------------------------------
# TestRenderPng
# ---------------------------------------------------------------------------

class TestRenderPng:
    def test_graceful_when_playwright_missing(self, tmp_path):
        html_file = tmp_path / "test.html"
        html_file.write_text("<html><body>test</body></html>")
        with patch.dict("sys.modules", {"playwright": None, "playwright.sync_api": None}):
            result = render_png(str(html_file), str(tmp_path / "out.png"))
            assert result is None


# ---------------------------------------------------------------------------
# TestStep5Integration
# ---------------------------------------------------------------------------

class TestStep5Integration:
    @patch("local_notebooklm.steps.step5.generate_text")
    @patch("local_notebooklm.steps.step5.render_png", return_value=None)
    def test_full_pipeline(self, mock_png, mock_gen, tmp_path):
        # Set up input
        input_dir = tmp_path / "step3"
        input_dir.mkdir()
        (input_dir / "podcast_ready_data.txt").write_text(
            "Speaker 1: Hello\nSpeaker 2: Hi there"
        )
        output_dir = tmp_path / "step5"

        mock_gen.return_value = json.dumps(SAMPLE_DATA)

        config = {
            "Big-Text-Model": {"model": "test-model"},
            "Step5": {"max_tokens": 4096, "temperature": 0.4},
        }

        result = step5(MagicMock(), config, str(input_dir), str(output_dir))

        assert result.endswith("infographic.html")
        assert (output_dir / "infographic.html").exists()
        assert (output_dir / "infographic_data.json").exists()

        # Verify JSON is valid
        with open(output_dir / "infographic_data.json") as f:
            saved = json.load(f)
        assert saved["title"] == "AI Revolution"

    def test_missing_transcript_raises(self, tmp_path):
        with pytest.raises(InfographicError, match="No transcript found"):
            step5(MagicMock(), {"Big-Text-Model": {"model": "m"}}, str(tmp_path), str(tmp_path / "out"))


# ---------------------------------------------------------------------------
# TestPptxGeneration
# ---------------------------------------------------------------------------

class TestPptxGeneration:
    def test_pptx_creation(self, tmp_path):
        out = tmp_path / "infographic.pptx"
        result = render_infographic_pptx(SAMPLE_DATA, str(out))
        assert result is not None
        assert out.exists()
        assert out.stat().st_size > 0

    def test_pptx_slide_count(self, tmp_path):
        from pptx import Presentation

        out = tmp_path / "infographic.pptx"
        render_infographic_pptx(SAMPLE_DATA, str(out))
        prs = Presentation(str(out))
        assert len(prs.slides) == 6

    def test_pptx_title_slide(self, tmp_path):
        from pptx import Presentation

        out = tmp_path / "infographic.pptx"
        render_infographic_pptx(SAMPLE_DATA, str(out))
        prs = Presentation(str(out))
        slide = prs.slides[0]
        texts = [
            shape.text_frame.paragraphs[0].text
            for shape in slide.shapes
            if shape.has_text_frame
        ]
        assert "AI Revolution" in texts

    def test_pptx_graceful_without_pptx(self, tmp_path):
        out = tmp_path / "infographic.pptx"
        with patch.dict("sys.modules", {"pptx": None, "pptx.util": None,
                                         "pptx.enum.text": None}):
            result = render_infographic_pptx(SAMPLE_DATA, str(out))
            assert result is None

    def test_video_graceful_without_moviepy(self, tmp_path):
        out_pptx = tmp_path / "infographic.pptx"
        render_infographic_pptx(SAMPLE_DATA, str(out_pptx))
        out_mp4 = tmp_path / "infographic.mp4"
        with patch.dict("sys.modules", {"moviepy": None,
                                         "moviepy.video": None,
                                         "moviepy.video.io": None,
                                         "moviepy.video.io.ImageSequenceClip": None}):
            result = render_video(str(out_pptx), str(out_mp4))
            assert result is None

    @patch("local_notebooklm.steps.step5.generate_text")
    @patch("local_notebooklm.steps.step5.render_png", return_value=None)
    def test_step5_produces_pptx(self, mock_png, mock_gen, tmp_path):
        input_dir = tmp_path / "step3"
        input_dir.mkdir()
        (input_dir / "podcast_ready_data.txt").write_text(
            "Speaker 1: Hello\nSpeaker 2: Hi there"
        )
        output_dir = tmp_path / "step5"

        mock_gen.return_value = json.dumps(SAMPLE_DATA)

        config = {
            "Big-Text-Model": {"model": "test-model"},
            "Step5": {"max_tokens": 4096, "temperature": 0.4},
        }

        step5(MagicMock(), config, str(input_dir), str(output_dir))

        assert (output_dir / "infographic.html").exists()
        assert (output_dir / "infographic.pptx").exists()


# ---------------------------------------------------------------------------
# TestChartGeneration
# ---------------------------------------------------------------------------

class TestChartGeneration:
    def test_topic_chart_returns_base64(self):
        result = generate_topic_importance_chart(SAMPLE_DATA)
        assert result is not None
        b64, png = result
        assert b64.startswith("data:image/png;base64,")
        assert png is None

    def test_topic_chart_saves_png(self, tmp_path):
        out = str(tmp_path / "topics.png")
        result = generate_topic_importance_chart(SAMPLE_DATA, output_path=out)
        assert result is not None
        _, png = result
        assert png == out
        assert (tmp_path / "topics.png").exists()

    def test_topic_chart_empty_data(self):
        result = generate_topic_importance_chart({"topics": []})
        assert result is None

    def test_speaker_chart_returns_base64(self):
        result = generate_speaker_distribution_chart(SAMPLE_DATA)
        assert result is not None
        b64, png = result
        assert b64.startswith("data:image/png;base64,")
        assert png is None

    def test_speaker_chart_saves_png(self, tmp_path):
        out = str(tmp_path / "speakers.png")
        result = generate_speaker_distribution_chart(SAMPLE_DATA, output_path=out)
        assert result is not None
        _, png = result
        assert png == out
        assert (tmp_path / "speakers.png").exists()

    def test_speaker_chart_empty_data(self):
        result = generate_speaker_distribution_chart({"speakers": []})
        assert result is None

    def test_speaker_chart_zero_counts(self):
        data = {"speakers": [{"label": "A", "line_count": 0}]}
        result = generate_speaker_distribution_chart(data)
        assert result is None

    def test_flow_chart_returns_base64(self):
        result = generate_conversation_flow_chart(SAMPLE_DATA)
        assert result is not None
        b64, png = result
        assert b64.startswith("data:image/png;base64,")
        assert png is None

    def test_flow_chart_saves_png(self, tmp_path):
        out = str(tmp_path / "flow.png")
        result = generate_conversation_flow_chart(SAMPLE_DATA, output_path=out)
        assert result is not None
        _, png = result
        assert png == out
        assert (tmp_path / "flow.png").exists()

    def test_flow_chart_empty_data(self):
        result = generate_conversation_flow_chart({"conversation_flow": []})
        assert result is None

    def test_generate_all_charts(self, tmp_path):
        results = generate_all_charts(SAMPLE_DATA, str(tmp_path))
        assert "topics" in results
        assert "speakers" in results
        assert "flow" in results
        generated = sum(1 for v in results.values() if v is not None)
        assert generated == 3
        assert (tmp_path / "chart_topics.png").exists()
        assert (tmp_path / "chart_speakers.png").exists()
        assert (tmp_path / "chart_flow.png").exists()


# ---------------------------------------------------------------------------
# TestChartGracefulDegradation
# ---------------------------------------------------------------------------

class TestChartGracefulDegradation:
    def test_topic_chart_none_without_matplotlib(self):
        with patch.dict("sys.modules", {"matplotlib": None, "matplotlib.pyplot": None}):
            with patch("local_notebooklm.steps.step5_charts._setup_cyberpunk_style", side_effect=ImportError):
                result = generate_topic_importance_chart(SAMPLE_DATA)
                assert result is None

    def test_speaker_chart_none_without_matplotlib(self):
        with patch("local_notebooklm.steps.step5_charts._setup_cyberpunk_style", side_effect=ImportError):
            result = generate_speaker_distribution_chart(SAMPLE_DATA)
            assert result is None

    def test_flow_chart_none_without_matplotlib(self):
        with patch("local_notebooklm.steps.step5_charts._setup_cyberpunk_style", side_effect=ImportError):
            result = generate_conversation_flow_chart(SAMPLE_DATA)
            assert result is None

    def test_generate_all_returns_nones_without_matplotlib(self):
        with patch("local_notebooklm.steps.step5_charts._setup_cyberpunk_style", side_effect=ImportError):
            results = generate_all_charts(SAMPLE_DATA)
            assert all(v is None for v in results.values())


# ---------------------------------------------------------------------------
# TestHtmlWithCharts
# ---------------------------------------------------------------------------

class TestHtmlWithCharts:
    def test_no_charts_no_img_tags(self):
        result = render_infographic_html(SAMPLE_DATA)
        assert "data:image/png;base64," not in result

    def test_with_charts_has_img_tags(self):
        fake_b64 = "data:image/png;base64,AAAA"
        charts = {
            "topics": (fake_b64, None),
            "speakers": (fake_b64, None),
            "flow": (fake_b64, None),
        }
        result = render_infographic_html(SAMPLE_DATA, charts=charts)
        assert result.count("data:image/png;base64,AAAA") == 3
        assert result.count("<img ") == 3

    def test_partial_charts_only_available(self):
        fake_b64 = "data:image/png;base64,BBBB"
        charts = {
            "topics": (fake_b64, None),
            "speakers": None,
            "flow": None,
        }
        result = render_infographic_html(SAMPLE_DATA, charts=charts)
        assert result.count("data:image/png;base64,BBBB") == 1
        assert result.count("<img ") == 1

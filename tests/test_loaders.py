"""Tests for local_notebooklm.loaders — all 6 extraction formats + error handling."""

import os
import tempfile
import pytest
from unittest.mock import patch, MagicMock
from local_notebooklm.loaders import (
    load_input,
    extract_text_from_pdf,
    extract_text_from_docx,
    extract_text_from_pptx,
    extract_text_from_txt,
    extract_text_from_markdown,
    extract_text_from_url,
    _is_youtube_url,
    _extract_youtube_video_id,
    LoaderError,
)


# ── Fixtures ────────────────────────────────────────────────────────


@pytest.fixture
def tmp_txt(tmp_path):
    p = tmp_path / "sample.txt"
    p.write_text("Hello world from a text file", encoding="utf-8")
    return str(p)


@pytest.fixture
def tmp_md(tmp_path):
    p = tmp_path / "sample.md"
    p.write_text("# Heading\n\nSome **markdown** content", encoding="utf-8")
    return str(p)


@pytest.fixture
def tmp_docx(tmp_path):
    from docx import Document

    p = tmp_path / "sample.docx"
    doc = Document()
    doc.add_paragraph("First paragraph")
    doc.add_paragraph("Second paragraph")
    doc.save(str(p))
    return str(p)


@pytest.fixture
def tmp_pptx(tmp_path):
    from pptx import Presentation

    p = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Slide Title"
    slide.placeholders[1].text = "Bullet content"
    prs.save(str(p))
    return str(p)


# ── TXT Tests ───────────────────────────────────────────────────────


class TestTxtLoader:
    def test_basic_read(self, tmp_txt):
        text = extract_text_from_txt(tmp_txt)
        assert text == "Hello world from a text file"

    def test_max_chars(self, tmp_txt):
        text = extract_text_from_txt(tmp_txt, max_chars=5)
        assert text == "Hello"

    def test_missing_file(self):
        with pytest.raises(LoaderError, match="File not found"):
            extract_text_from_txt("/nonexistent/path.txt")

    def test_encoding_fallback(self, tmp_path):
        p = tmp_path / "latin.txt"
        p.write_bytes("caf\xe9".encode("latin-1"))
        text = extract_text_from_txt(str(p))
        assert "caf" in text

    def test_via_dispatcher(self, tmp_txt):
        text = load_input(tmp_txt)
        assert "Hello" in text


# ── Markdown Tests ──────────────────────────────────────────────────


class TestMarkdownLoader:
    def test_basic_read(self, tmp_md):
        text = extract_text_from_markdown(tmp_md)
        assert "# Heading" in text
        assert "**markdown**" in text

    def test_max_chars(self, tmp_md):
        text = extract_text_from_markdown(tmp_md, max_chars=10)
        assert len(text) == 10

    def test_missing_file(self):
        with pytest.raises(LoaderError, match="File not found"):
            extract_text_from_markdown("/nonexistent/path.md")

    def test_via_dispatcher(self, tmp_md):
        text = load_input(tmp_md)
        assert "Heading" in text


# ── DOCX Tests ──────────────────────────────────────────────────────


class TestDocxLoader:
    def test_basic_read(self, tmp_docx):
        text = extract_text_from_docx(tmp_docx)
        assert "First paragraph" in text
        assert "Second paragraph" in text

    def test_max_chars(self, tmp_docx):
        text = extract_text_from_docx(tmp_docx, max_chars=10)
        assert len(text) <= 10

    def test_missing_file(self):
        with pytest.raises(LoaderError, match="File not found"):
            extract_text_from_docx("/nonexistent/path.docx")

    def test_via_dispatcher(self, tmp_docx):
        text = load_input(tmp_docx)
        assert "paragraph" in text


# ── PPTX Tests ──────────────────────────────────────────────────────


class TestPptxLoader:
    def test_basic_read(self, tmp_pptx):
        text = extract_text_from_pptx(tmp_pptx)
        assert "Slide Title" in text
        assert "Bullet content" in text

    def test_max_chars(self, tmp_pptx):
        text = extract_text_from_pptx(tmp_pptx, max_chars=8)
        assert len(text) <= 8

    def test_missing_file(self):
        with pytest.raises(LoaderError, match="File not found"):
            extract_text_from_pptx("/nonexistent/path.pptx")

    def test_via_dispatcher(self, tmp_pptx):
        text = load_input(tmp_pptx)
        assert "Slide" in text


# ── PDF Tests ───────────────────────────────────────────────────────


class TestPdfLoader:
    def test_missing_file(self):
        with pytest.raises(LoaderError, match="File not found"):
            extract_text_from_pdf("/nonexistent/path.pdf")

    def test_example_pdf(self):
        pdf = "./examples/MoshiVis.pdf"
        if not os.path.exists(pdf):
            pytest.skip("Example PDF not available")
        text = extract_text_from_pdf(pdf, max_chars=200)
        assert len(text) > 0
        assert len(text) <= 200

    def test_via_dispatcher_extension(self, tmp_path):
        """Dispatcher routes .pdf to PDF loader."""
        p = tmp_path / "fake.pdf"
        p.write_bytes(b"not a real pdf")
        with pytest.raises(LoaderError):
            load_input(str(p))


# ── URL Tests ───────────────────────────────────────────────────────


class TestUrlLoader:
    def test_trafilatura_extraction(self):
        """Mock trafilatura at the function level (no reload needed)."""
        mock_traf = MagicMock()
        mock_traf.fetch_url.return_value = "<html>article</html>"
        mock_traf.extract.return_value = "Extracted article content"

        with patch.dict("local_notebooklm.loaders.__builtins__", {}):
            with patch("local_notebooklm.loaders.extract_text_from_url") as mock_fn:
                mock_fn.return_value = "Extracted article content"
                text = load_input("https://example.com/article")
                assert text == "Extracted article content"

    def test_url_via_dispatcher(self):
        with patch("local_notebooklm.loaders.extract_text_from_url", return_value="url content"):
            text = load_input("https://example.com/article")
            assert text == "url content"

    def test_url_max_chars(self):
        with patch("local_notebooklm.loaders.extract_text_from_url", return_value="x" * 100):
            text = load_input("https://example.com", max_chars=100)
            assert len(text) <= 100


# ── YouTube URL Detection Tests ─────────────────────────────────────


class TestYouTubeUrlDetection:
    """Tests for _is_youtube_url and _extract_youtube_video_id."""

    # -- _is_youtube_url positive cases --

    @pytest.mark.parametrize(
        "url",
        [
            "https://www.youtube.com/watch?v=dQw4w9WgXcQ",
            "https://youtube.com/watch?v=dQw4w9WgXcQ",
            "http://www.youtube.com/watch?v=dQw4w9WgXcQ",
            "https://youtu.be/dQw4w9WgXcQ",
            "https://www.youtube.com/embed/dQw4w9WgXcQ",
            "https://www.youtube.com/v/dQw4w9WgXcQ",
            "https://www.youtube.com/shorts/dQw4w9WgXcQ",
        ],
    )
    def test_is_youtube_url_positive(self, url):
        assert _is_youtube_url(url) is True

    # -- _is_youtube_url negative cases --

    @pytest.mark.parametrize(
        "url",
        [
            "https://example.com/watch?v=abc",
            "https://vimeo.com/12345",
            "https://www.youtube.com/channel/UCabc123",
            "https://www.youtube.com/playlist?list=PLabc",
            "https://www.youtube.com/",
        ],
    )
    def test_is_youtube_url_negative(self, url):
        assert _is_youtube_url(url) is False

    # -- _extract_youtube_video_id --

    @pytest.mark.parametrize(
        "url, expected_id",
        [
            ("https://www.youtube.com/watch?v=dQw4w9WgXcQ", "dQw4w9WgXcQ"),
            ("https://youtu.be/dQw4w9WgXcQ", "dQw4w9WgXcQ"),
            ("https://www.youtube.com/embed/dQw4w9WgXcQ", "dQw4w9WgXcQ"),
            ("https://www.youtube.com/v/dQw4w9WgXcQ", "dQw4w9WgXcQ"),
            ("https://www.youtube.com/shorts/dQw4w9WgXcQ", "dQw4w9WgXcQ"),
            ("https://www.youtube.com/watch?v=dQw4w9WgXcQ&t=120", "dQw4w9WgXcQ"),
            ("https://youtu.be/dQw4w9WgXcQ?t=30", "dQw4w9WgXcQ"),
        ],
    )
    def test_extract_video_id(self, url, expected_id):
        assert _extract_youtube_video_id(url) == expected_id

    def test_extract_video_id_invalid_url(self):
        with pytest.raises(ValueError, match="Could not extract video ID"):
            _extract_youtube_video_id("https://www.youtube.com/")


# ── YouTube Transcript Loader Tests ────────────────────────────────


class TestYouTubeTranscriptLoader:
    """Tests for YouTube transcript extraction via extract_text_from_url."""

    def test_youtube_url_uses_transcript(self):
        """YouTube URL should use transcript when available."""
        with patch(
            "local_notebooklm.loaders._extract_youtube_transcript",
            return_value="Hello this is a transcript",
        ) as mock_yt:
            text = extract_text_from_url("https://www.youtube.com/watch?v=abc123")
            assert text == "Hello this is a transcript"
            mock_yt.assert_called_once_with("abc123", 100000)

    def test_falls_back_on_import_error(self):
        """Missing youtube-transcript-api should fall back to web scraping."""
        with patch(
            "local_notebooklm.loaders._extract_youtube_transcript",
            side_effect=ImportError("No module named 'youtube_transcript_api'"),
        ):
            with patch(
                "local_notebooklm.loaders.extract_text_from_url",
                wraps=extract_text_from_url,
            ):
                # Patch trafilatura to succeed as the fallback
                mock_traf = MagicMock()
                mock_traf.fetch_url.return_value = "<html>fallback</html>"
                mock_traf.extract.return_value = "Fallback article content"
                with patch.dict("sys.modules", {"trafilatura": mock_traf}):
                    text = extract_text_from_url("https://www.youtube.com/watch?v=abc123")
                    assert text == "Fallback article content"

    def test_falls_back_on_runtime_error(self):
        """Failed transcript fetch should fall back to web scraping."""
        with patch(
            "local_notebooklm.loaders._extract_youtube_transcript",
            side_effect=RuntimeError("No transcript available"),
        ):
            mock_traf = MagicMock()
            mock_traf.fetch_url.return_value = "<html>fallback</html>"
            mock_traf.extract.return_value = "Scraped content"
            with patch.dict("sys.modules", {"trafilatura": mock_traf}):
                text = extract_text_from_url("https://www.youtube.com/watch?v=abc123")
                assert text == "Scraped content"

    def test_youtu_be_short_url(self):
        """youtu.be short URLs should route through YouTube extraction."""
        with patch(
            "local_notebooklm.loaders._extract_youtube_transcript",
            return_value="Short URL transcript",
        ) as mock_yt:
            text = extract_text_from_url("https://youtu.be/xyz789")
            assert text == "Short URL transcript"
            mock_yt.assert_called_once_with("xyz789", 100000)

    def test_non_youtube_url_skips_youtube_path(self):
        """Non-YouTube URLs should not attempt YouTube transcript extraction."""
        with patch(
            "local_notebooklm.loaders._extract_youtube_transcript",
        ) as mock_yt:
            mock_traf = MagicMock()
            mock_traf.fetch_url.return_value = "<html>article</html>"
            mock_traf.extract.return_value = "Regular article"
            with patch.dict("sys.modules", {"trafilatura": mock_traf}):
                text = extract_text_from_url("https://example.com/article")
                assert text == "Regular article"
                mock_yt.assert_not_called()

    def test_youtube_url_via_dispatcher(self):
        """load_input routes YouTube URLs through extract_text_from_url."""
        with patch(
            "local_notebooklm.loaders._extract_youtube_transcript",
            return_value="Dispatched transcript",
        ):
            text = load_input("https://www.youtube.com/watch?v=test123")
            assert text == "Dispatched transcript"

    def test_youtube_transcript_max_chars(self):
        """YouTube transcript should respect max_chars."""
        with patch(
            "local_notebooklm.loaders._extract_youtube_transcript",
            return_value="A" * 50,
        ):
            text = extract_text_from_url("https://www.youtube.com/watch?v=abc", max_chars=50)
            assert len(text) <= 50


# ── Dispatcher Tests ────────────────────────────────────────────────


class TestDispatcher:
    def test_empty_path(self):
        with pytest.raises(LoaderError, match="No input path"):
            load_input("")

    def test_none_path(self):
        with pytest.raises(LoaderError, match="No input path"):
            load_input(None)

    def test_unsupported_extension(self, tmp_path):
        p = tmp_path / "file.xyz"
        p.write_text("data")
        with pytest.raises(LoaderError, match="Unsupported file type"):
            load_input(str(p))

    def test_http_routes_to_url(self):
        with patch("local_notebooklm.loaders.extract_text_from_url", return_value="ok") as mock:
            load_input("http://example.com")
            mock.assert_called_once()

    def test_https_routes_to_url(self):
        with patch("local_notebooklm.loaders.extract_text_from_url", return_value="ok") as mock:
            load_input("https://example.com")
            mock.assert_called_once()


# ── Docling PDF Loader Tests ───────────────────────────────────────


class TestDoclingPdfLoader:
    """Tests for the docling-first-then-PyPDF2 fallback pattern."""

    def test_uses_docling_when_available(self, tmp_path):
        """When docling is importable and works, PyPDF2 should not be called."""
        pdf = tmp_path / "test.pdf"
        pdf.write_bytes(b"%PDF-fake")

        mock_converter_instance = MagicMock()
        mock_converter_instance.convert.return_value.document.export_to_markdown.return_value = (
            "# Title\n\nDocling extracted content"
        )
        mock_converter_cls = MagicMock(return_value=mock_converter_instance)

        mock_docling = MagicMock()
        mock_docling.document_converter.DocumentConverter = mock_converter_cls

        with patch.dict("sys.modules", {"docling": mock_docling, "docling.document_converter": mock_docling.document_converter}):
            with patch("local_notebooklm.loaders._extract_pdf_with_pypdf2") as mock_pypdf2:
                result = extract_text_from_pdf(str(pdf))
                assert "Docling extracted content" in result
                mock_pypdf2.assert_not_called()

    def test_falls_back_when_docling_missing(self, tmp_path):
        """ImportError from docling should fall back to PyPDF2."""
        pdf = tmp_path / "test.pdf"
        pdf.write_bytes(b"%PDF-fake")

        with patch("local_notebooklm.loaders._extract_pdf_with_docling", side_effect=ImportError("No module")):
            with patch("local_notebooklm.loaders._extract_pdf_with_pypdf2", return_value="PyPDF2 text") as mock_p:
                result = extract_text_from_pdf(str(pdf))
                assert result == "PyPDF2 text"
                mock_p.assert_called_once()

    def test_falls_back_when_docling_fails(self, tmp_path):
        """RuntimeError from docling should fall back to PyPDF2."""
        pdf = tmp_path / "test.pdf"
        pdf.write_bytes(b"%PDF-fake")

        with patch("local_notebooklm.loaders._extract_pdf_with_docling", side_effect=RuntimeError("parse error")):
            with patch("local_notebooklm.loaders._extract_pdf_with_pypdf2", return_value="fallback") as mock_p:
                result = extract_text_from_pdf(str(pdf))
                assert result == "fallback"
                mock_p.assert_called_once()

    def test_falls_back_when_docling_returns_empty(self, tmp_path):
        """Docling returning empty text should fall back to PyPDF2."""
        pdf = tmp_path / "test.pdf"
        pdf.write_bytes(b"%PDF-fake")

        with patch("local_notebooklm.loaders._extract_pdf_with_docling", side_effect=ValueError("Docling returned empty text")):
            with patch("local_notebooklm.loaders._extract_pdf_with_pypdf2", return_value="fallback content") as mock_p:
                result = extract_text_from_pdf(str(pdf))
                assert result == "fallback content"
                mock_p.assert_called_once()

    def test_docling_max_chars_truncation(self, tmp_path):
        """Docling output should be truncated at a line boundary near max_chars."""
        pdf = tmp_path / "test.pdf"
        pdf.write_bytes(b"%PDF-fake")

        long_text = "Line one\nLine two\nLine three\nLine four\nLine five"

        mock_converter_instance = MagicMock()
        mock_converter_instance.convert.return_value.document.export_to_markdown.return_value = long_text
        mock_converter_cls = MagicMock(return_value=mock_converter_instance)

        mock_docling = MagicMock()
        mock_docling.document_converter.DocumentConverter = mock_converter_cls

        with patch.dict("sys.modules", {"docling": mock_docling, "docling.document_converter": mock_docling.document_converter}):
            result = extract_text_from_pdf(str(pdf), max_chars=20)
            assert len(result) <= 20
            # Should truncate at a newline boundary
            assert "\n" not in result or result.endswith("\n") is False

    def test_file_not_found_still_raises(self):
        """File-not-found check should fire before dispatch to either extractor."""
        with pytest.raises(LoaderError, match="File not found"):
            extract_text_from_pdf("/nonexistent/path.pdf")

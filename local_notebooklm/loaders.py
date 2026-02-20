"""Multi-format document loaders for Local-NotebookLM.

Supports PDF, DOCX, PPTX, TXT, Markdown, and web URLs.
"""

import logging
import os
import re
from pathlib import Path
from urllib.parse import urlparse, parse_qs

logger = logging.getLogger(__name__)


class LoaderError(Exception):
    pass


def _is_youtube_url(url: str) -> bool:
    """Return True if *url* points to a YouTube video."""
    return bool(
        re.match(
            r"https?://(?:www\.|m\.)?(?:youtube\.com/(?:watch|embed|v|shorts)/?"
            r"|youtu\.be/)",
            url,
        )
    )


def _extract_youtube_video_id(url: str) -> str:
    """Parse a YouTube video ID from any common URL format."""
    parsed = urlparse(url)
    host = parsed.hostname or ""

    # youtu.be/<id>
    if host in ("youtu.be", "www.youtu.be"):
        video_id = parsed.path.lstrip("/").split("/")[0]
        if video_id:
            return video_id

    # youtube.com/watch?v=<id>
    if "v" in parse_qs(parsed.query):
        return parse_qs(parsed.query)["v"][0]

    # youtube.com/embed/<id>, /v/<id>, /shorts/<id>
    parts = [p for p in parsed.path.split("/") if p]
    if len(parts) >= 2 and parts[0] in ("embed", "v", "shorts"):
        return parts[1]

    raise ValueError(f"Could not extract video ID from URL: {url}")


def _extract_youtube_transcript(video_id: str, max_chars: int = 100000) -> str:
    """Fetch a YouTube transcript using youtube-transcript-api."""
    from youtube_transcript_api import YouTubeTranscriptApi

    ytt = YouTubeTranscriptApi()
    try:
        transcript = ytt.fetch(video_id, languages=["en"])
    except Exception:
        # Fall back to any available language
        transcript = ytt.fetch(video_id)

    text = " ".join(snippet.text for snippet in transcript)
    text = text[:max_chars]
    logger.info(f"YouTube transcript extraction complete. {len(text)} chars")
    return text


def _extract_pdf_with_docling(file_path: str, max_chars: int) -> str:
    """Extract PDF text using docling (higher-quality: tables, layout, OCR)."""
    from docling.document_converter import DocumentConverter

    converter = DocumentConverter()
    result = converter.convert(file_path)
    text = result.document.export_to_markdown()

    if not text or not text.strip():
        raise ValueError("Docling returned empty text")

    # Truncate at a line boundary near max_chars
    if len(text) > max_chars:
        cut = text.rfind("\n", 0, max_chars)
        if cut == -1:
            cut = max_chars
        text = text[:cut]

    logger.info(f"PDF extraction (docling) complete. {len(text)} chars")
    return text


def _extract_pdf_with_pypdf2(file_path: str, max_chars: int) -> str:
    """Extract PDF text using PyPDF2 (fallback)."""
    import PyPDF2

    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        num_pages = len(reader.pages)
        logger.info(f"Processing PDF with {num_pages} pages (PyPDF2)")

        parts = []
        total = 0
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if total + len(text) > max_chars:
                parts.append(text[: max_chars - total])
                logger.info(f"Reached {max_chars} char limit at page {i + 1}")
                break
            parts.append(text)
            total += len(text)

        result = "\n".join(parts)
        logger.info(f"PDF extraction (PyPDF2) complete. {len(result)} chars")
        return result


def extract_text_from_pdf(file_path: str, max_chars: int = 100000) -> str:
    if not os.path.exists(file_path):
        raise LoaderError(f"File not found: {file_path}")

    # Try docling first (higher-quality extraction with tables/layout/OCR)
    try:
        return _extract_pdf_with_docling(file_path, max_chars)
    except ImportError:
        logger.info("docling not installed, falling back to PyPDF2")
    except Exception as e:
        logger.warning(f"docling extraction failed, falling back to PyPDF2: {e}")

    # Fallback to PyPDF2
    try:
        return _extract_pdf_with_pypdf2(file_path, max_chars)
    except Exception as e:
        raise LoaderError(f"Failed to extract text from PDF: {e}")


def extract_text_from_docx(file_path: str, max_chars: int = 100000) -> str:
    from docx import Document

    if not os.path.exists(file_path):
        raise LoaderError(f"File not found: {file_path}")

    try:
        doc = Document(file_path)
        parts = []
        total = 0

        def _add(text: str) -> bool:
            nonlocal total
            if not text:
                return True
            if total + len(text) + 1 > max_chars:
                parts.append(text[: max_chars - total])
                return False
            parts.append(text)
            total += len(text) + 1
            return True

        for para in doc.paragraphs:
            if not _add(para.text):
                break

        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells and not _add(" | ".join(cells)):
                    break

        result = "\n".join(parts)
        logger.info(f"DOCX extraction complete. {len(result)} chars")
        return result
    except LoaderError:
        raise
    except Exception as e:
        raise LoaderError(f"Failed to extract text from DOCX: {e}")


def extract_text_from_pptx(file_path: str, max_chars: int = 100000) -> str:
    from pptx import Presentation

    if not os.path.exists(file_path):
        raise LoaderError(f"File not found: {file_path}")

    try:
        prs = Presentation(file_path)
        parts = []
        total = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue
                    if total + len(text) + 1 > max_chars:
                        parts.append(text[: max_chars - total])
                        result = "\n".join(parts)
                        logger.info(f"PPTX extraction complete. {len(result)} chars")
                        return result
                    parts.append(text)
                    total += len(text) + 1

        result = "\n".join(parts)
        logger.info(f"PPTX extraction complete. {len(result)} chars")
        return result
    except LoaderError:
        raise
    except Exception as e:
        raise LoaderError(f"Failed to extract text from PPTX: {e}")


def extract_text_from_txt(file_path: str, max_chars: int = 100000) -> str:
    if not os.path.exists(file_path):
        raise LoaderError(f"File not found: {file_path}")

    for encoding in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
        try:
            with open(file_path, "r", encoding=encoding) as f:
                text = f.read(max_chars)
            logger.info(f"TXT extraction complete ({encoding}). {len(text)} chars")
            return text
        except UnicodeDecodeError:
            continue

    raise LoaderError(f"Could not decode {file_path} with any supported encoding")


def extract_text_from_markdown(file_path: str, max_chars: int = 100000) -> str:
    if not os.path.exists(file_path):
        raise LoaderError(f"File not found: {file_path}")

    try:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read(max_chars)
        logger.info(f"Markdown extraction complete. {len(text)} chars")
        return text
    except Exception as e:
        raise LoaderError(f"Failed to read Markdown file: {e}")


def extract_text_from_url(url: str, max_chars: int = 100000) -> str:
    # YouTube detection — extract transcript instead of scraping HTML
    if _is_youtube_url(url):
        try:
            video_id = _extract_youtube_video_id(url)
            return _extract_youtube_transcript(video_id, max_chars)
        except ImportError:
            logger.info("youtube-transcript-api not installed, falling back to web scraping")
        except Exception as e:
            logger.warning(f"YouTube transcript extraction failed, falling back to web scraping: {e}")

    # Try trafilatura first (best article extraction)
    try:
        import trafilatura

        downloaded = trafilatura.fetch_url(url)
        if downloaded:
            text = trafilatura.extract(downloaded)
            if text:
                text = text[:max_chars]
                logger.info(f"URL extraction (trafilatura) complete. {len(text)} chars")
                return text
    except Exception as e:
        logger.warning(f"trafilatura failed, falling back to bs4: {e}")

    # Fallback to requests + BeautifulSoup
    try:
        import requests
        from bs4 import BeautifulSoup

        resp = requests.get(url, timeout=30, headers={"User-Agent": "Local-NotebookLM/1.0"})
        resp.raise_for_status()
        soup = BeautifulSoup(resp.text, "html.parser")

        # Remove script/style tags
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()

        text = soup.get_text(separator="\n", strip=True)
        text = text[:max_chars]
        logger.info(f"URL extraction (bs4) complete. {len(text)} chars")
        return text
    except Exception as e:
        raise LoaderError(f"Failed to extract text from URL: {e}")


# Extension-to-loader mapping
_LOADERS = {
    ".pdf": extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".pptx": extract_text_from_pptx,
    ".txt": extract_text_from_txt,
    ".md": extract_text_from_markdown,
}


def load_input(input_path: str, max_chars: int = 100000) -> str:
    """Dispatch to the correct loader based on file extension or URL prefix."""
    if not input_path:
        raise LoaderError("No input path provided")

    # URL detection
    if input_path.startswith(("http://", "https://")):
        return extract_text_from_url(input_path, max_chars)

    # File-based detection
    ext = Path(input_path).suffix.lower()
    loader = _LOADERS.get(ext)
    if loader is None:
        supported = ", ".join(sorted(_LOADERS.keys()))
        raise LoaderError(
            f"Unsupported file type '{ext}'. Supported: {supported} and URLs (http/https)"
        )

    return loader(input_path, max_chars)

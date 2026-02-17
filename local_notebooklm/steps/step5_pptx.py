"""Step 5d/5e — PPTX slide deck and optional video generation.

Generates a 6-slide PPTX presentation from structured infographic data,
with an optional MP4 video conversion.  Both outputs are non-fatal —
if python-pptx or moviepy are missing, the functions return None.
"""

import logging
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Theme configuration (cyberpunk dark, matching HTML infographic)
# ---------------------------------------------------------------------------

_THEME = {
    "bg": "161c24",        # dark blue-black
    "title": "00ffff",     # neon cyan
    "accent": "40c4ff",    # bright cyan
    "body": "e0e0f0",      # light gray
    "muted": "8888aa",     # muted gray
    "magenta": "dd368a",   # pink accent
}

# Slide dimensions: 16:9 widescreen
_SLIDE_WIDTH_EMU = 12192000   # 13.33 inches
_SLIDE_HEIGHT_EMU = 6858000   # 7.5 inches


def _hex_to_rgbcolor(hex_str: str):
    """Convert a hex string like '00ffff' to pptx RGBColor."""
    try:
        from pptx.dml.color import RGBColor
    except ImportError:
        from pptx.util import RGBColor
    return RGBColor.from_string(hex_str)


def _set_slide_bg(slide, hex_color: str):
    """Set a solid background color on a slide."""
    from pptx.util import Emu
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgbcolor(hex_color)


def _add_text_box(slide, left, top, width, height, text: str,
                  font_size: int = 14, color: str = "e0e0f0",
                  bold: bool = False, alignment=None):
    """Add a text box to a slide with styled text."""
    from pptx.util import Inches, Pt, Emu

    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = _hex_to_rgbcolor(color)
    p.font.bold = bold
    if alignment is not None:
        p.alignment = alignment
    return txBox


def _add_bullet_frame(slide, left, top, width, height, items: List[str],
                      font_size: int = 12, color: str = "e0e0f0",
                      prefix_fn=None):
    """Add a text frame with bullet items."""
    from pptx.util import Pt, Emu

    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        prefix = prefix_fn(i) if prefix_fn else ""
        p.text = f"{prefix}{item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = _hex_to_rgbcolor(color)
        p.space_after = Pt(6)
    return txBox


# ---------------------------------------------------------------------------
# 5d — PPTX generation
# ---------------------------------------------------------------------------

def render_infographic_pptx(
    data: Dict[str, Any],
    output_path: str,
) -> Optional[str]:
    """Generate a 6-slide PPTX presentation from structured infographic data.

    Returns the output path on success, or None if python-pptx is not installed.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        logger.info("python-pptx not installed -- skipping PPTX generation.")
        return None

    try:
        prs = Presentation()
        prs.slide_width = Emu(_SLIDE_WIDTH_EMU)
        prs.slide_height = Emu(_SLIDE_HEIGHT_EMU)

        # Use blank layout for full control
        blank_layout = prs.slide_layouts[6]  # blank

        # ----- Slide 1: Title -----
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide, _THEME["bg"])

        title_text = data.get("title", "Podcast Infographic")
        _add_text_box(
            slide,
            left=_SLIDE_WIDTH_EMU // 8,
            top=_SLIDE_HEIGHT_EMU // 4,
            width=_SLIDE_WIDTH_EMU * 3 // 4,
            height=Emu(Inches(1.5)),
            text=title_text,
            font_size=36,
            color=_THEME["title"],
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        summary = data.get("summary", "")
        if summary:
            _add_text_box(
                slide,
                left=_SLIDE_WIDTH_EMU // 8,
                top=_SLIDE_HEIGHT_EMU // 2,
                width=_SLIDE_WIDTH_EMU * 3 // 4,
                height=Emu(Inches(2)),
                text=summary,
                font_size=16,
                color=_THEME["muted"],
                alignment=PP_ALIGN.CENTER,
            )

        # ----- Slide 2: Topics -----
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide, _THEME["bg"])
        _add_text_box(
            slide,
            left=Emu(Inches(0.5)),
            top=Emu(Inches(0.3)),
            width=Emu(Inches(12)),
            height=Emu(Inches(0.6)),
            text="TOPICS",
            font_size=24,
            color=_THEME["title"],
            bold=True,
        )

        topics = data.get("topics", [])
        topic_lines = []
        for t in topics:
            name = t.get("name", "")
            desc = t.get("description", "")
            importance = t.get("importance", 3)
            stars = "*" * int(importance)
            topic_lines.append(f"[{stars}] {name} -- {desc}")

        if topic_lines:
            _add_bullet_frame(
                slide,
                left=Emu(Inches(0.5)),
                top=Emu(Inches(1.2)),
                width=Emu(Inches(12)),
                height=Emu(Inches(5.5)),
                items=topic_lines,
                font_size=14,
                color=_THEME["body"],
            )

        # ----- Slide 3: Key Takeaways -----
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide, _THEME["bg"])
        _add_text_box(
            slide,
            left=Emu(Inches(0.5)),
            top=Emu(Inches(0.3)),
            width=Emu(Inches(12)),
            height=Emu(Inches(0.6)),
            text="KEY TAKEAWAYS",
            font_size=24,
            color=_THEME["title"],
            bold=True,
        )

        takeaways = data.get("key_takeaways", [])
        if takeaways:
            _add_bullet_frame(
                slide,
                left=Emu(Inches(0.5)),
                top=Emu(Inches(1.2)),
                width=Emu(Inches(12)),
                height=Emu(Inches(5.5)),
                items=takeaways,
                font_size=14,
                color=_THEME["body"],
                prefix_fn=lambda i: f"{i + 1}. ",
            )

        # ----- Slide 4: Notable Quotes -----
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide, _THEME["bg"])
        _add_text_box(
            slide,
            left=Emu(Inches(0.5)),
            top=Emu(Inches(0.3)),
            width=Emu(Inches(12)),
            height=Emu(Inches(0.6)),
            text="NOTABLE QUOTES",
            font_size=24,
            color=_THEME["title"],
            bold=True,
        )

        quotes = data.get("notable_quotes", [])
        quote_lines = []
        for q in quotes:
            speaker = q.get("speaker", "")
            quote_text = q.get("quote", "")
            quote_lines.append(f'"{quote_text}" -- {speaker}')

        if quote_lines:
            _add_bullet_frame(
                slide,
                left=Emu(Inches(0.5)),
                top=Emu(Inches(1.2)),
                width=Emu(Inches(12)),
                height=Emu(Inches(5.5)),
                items=quote_lines,
                font_size=14,
                color=_THEME["accent"],
            )

        # ----- Slide 5: Speakers -----
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide, _THEME["bg"])
        _add_text_box(
            slide,
            left=Emu(Inches(0.5)),
            top=Emu(Inches(0.3)),
            width=Emu(Inches(12)),
            height=Emu(Inches(0.6)),
            text="SPEAKERS",
            font_size=24,
            color=_THEME["title"],
            bold=True,
        )

        speakers = data.get("speakers", [])
        speaker_lines = []
        for s in speakers:
            label = s.get("label", "")
            role = s.get("role", "")
            line_count = s.get("line_count", "?")
            speaker_lines.append(f"{label} ({role}) -- {line_count} turns")

        if speaker_lines:
            _add_bullet_frame(
                slide,
                left=Emu(Inches(0.5)),
                top=Emu(Inches(1.2)),
                width=Emu(Inches(12)),
                height=Emu(Inches(5.5)),
                items=speaker_lines,
                font_size=16,
                color=_THEME["body"],
            )

        # ----- Slide 6: Conversation Flow -----
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide, _THEME["bg"])
        _add_text_box(
            slide,
            left=Emu(Inches(0.5)),
            top=Emu(Inches(0.3)),
            width=Emu(Inches(12)),
            height=Emu(Inches(0.6)),
            text="CONVERSATION FLOW",
            font_size=24,
            color=_THEME["title"],
            bold=True,
        )

        flow = data.get("conversation_flow", [])
        flow_lines = []
        for i, entry in enumerate(flow, 1):
            speaker = entry.get("speaker", "")
            topic = entry.get("topic", "")
            flow_lines.append(f"{speaker}: {topic}")

        if flow_lines:
            _add_bullet_frame(
                slide,
                left=Emu(Inches(0.5)),
                top=Emu(Inches(1.2)),
                width=Emu(Inches(12)),
                height=Emu(Inches(5.5)),
                items=flow_lines,
                font_size=14,
                color=_THEME["body"],
                prefix_fn=lambda i: f"{i + 1}. ",
            )

        # Save
        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out))
        logger.info(f"PPTX infographic saved to {output_path}")
        return str(out)

    except Exception as exc:
        logger.warning(f"PPTX generation failed: {exc}")
        return None


# ---------------------------------------------------------------------------
# 5e — Optional video conversion (PPTX slides → frames → MP4)
# ---------------------------------------------------------------------------

def render_video(
    pptx_path: str,
    output_path: str,
    seconds_per_slide: int = 3,
    fps: int = 24,
    resolution: tuple = (1920, 1080),
) -> Optional[str]:
    """Convert PPTX slides to MP4 video using PIL + moviepy.

    Renders each slide's text content onto a dark background image,
    then stitches frames into an MP4 with moviepy.

    Returns the output path on success, or None if dependencies are missing.
    """
    try:
        from pptx import Presentation
    except ImportError:
        logger.info("python-pptx not installed -- skipping video generation.")
        return None

    try:
        from moviepy.video.io.ImageSequenceClip import ImageSequenceClip
    except ImportError:
        logger.info("moviepy not installed -- skipping video generation.")
        return None

    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        logger.info("Pillow not installed -- skipping video generation.")
        return None

    try:
        prs = Presentation(pptx_path)
        width, height = resolution
        frames = []

        # Parse hex colors
        bg_rgb = tuple(int(_THEME["bg"][i:i+2], 16) for i in (0, 2, 4))
        title_rgb = tuple(int(_THEME["title"][i:i+2], 16) for i in (0, 2, 4))
        body_rgb = tuple(int(_THEME["body"][i:i+2], 16) for i in (0, 2, 4))

        for slide in prs.slides:
            img = Image.new("RGB", (width, height), bg_rgb)
            draw = ImageDraw.Draw(img)

            # Extract all text from the slide
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)

            # Draw title (first line) and body (remaining)
            y_pos = 80
            for i, text in enumerate(texts):
                color = title_rgb if i == 0 else body_rgb
                font_size_val = 48 if i == 0 else 28

                try:
                    font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", font_size_val)
                except (OSError, IOError):
                    font = ImageFont.load_default()

                # Word-wrap long lines
                max_chars = width // (font_size_val // 2 + 2)
                lines = []
                while text:
                    if len(text) <= max_chars:
                        lines.append(text)
                        break
                    split_at = text.rfind(" ", 0, max_chars)
                    if split_at == -1:
                        split_at = max_chars
                    lines.append(text[:split_at])
                    text = text[split_at:].lstrip()

                for line in lines:
                    if y_pos > height - 60:
                        break
                    draw.text((60, y_pos), line, fill=color, font=font)
                    y_pos += font_size_val + 8

                y_pos += 16  # spacing between items

            # Each slide becomes N frames (seconds_per_slide * fps)
            import tempfile
            frame_path = tempfile.mktemp(suffix=".png")
            img.save(frame_path)
            frames.extend([frame_path] * (seconds_per_slide * fps))

        if not frames:
            logger.warning("No frames generated from PPTX slides.")
            return None

        # Stitch into video
        clip = ImageSequenceClip(frames, fps=fps)
        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        clip.write_videofile(
            str(out),
            fps=fps,
            codec="libx264",
            audio=False,
            logger=None,
        )

        # Clean up temp frame files
        for f in set(frames):
            try:
                Path(f).unlink(missing_ok=True)
            except Exception:
                pass

        logger.info(f"Video saved to {output_path}")
        return str(out)

    except Exception as exc:
        logger.warning(f"Video generation failed: {exc}")
        return None
